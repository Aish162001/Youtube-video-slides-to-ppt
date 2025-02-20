import yt_dlp
import cv2
import os
from skimage.metrics import structural_similarity as ssim
from pptx import Presentation
from pptx.util import Inches
import re
from natsort import natsorted  # Import natural sorting

def download_video(url, output_path="video.mp4"):
    ydl_opts = {
        'format': 'best',
        'outtmpl': output_path,
    }
    with yt_dlp.YoutubeDL(ydl_opts) as ydl:
        ydl.download([url])
    return output_path

video_url = "VIDEO_URL"
video_path = download_video(video_url)
print(f"Video downloaded to: {video_path}")


def extract_frames(video_path, output_folder="frames", interval=1):
    os.makedirs(output_folder, exist_ok=True)
    cap = cv2.VideoCapture(video_path)
    frame_rate = int(cap.get(cv2.CAP_PROP_FPS))  # Frames per second
    count = 0
    saved_frames = []

    while cap.isOpened():
        ret, frame = cap.read()
        if not ret:
            break
        if count % (frame_rate * interval) == 0:  # Capture every interval second
            frame_path = os.path.join(output_folder, f"frame_{count}.jpg")
            cv2.imwrite(frame_path, frame)
            saved_frames.append(frame_path)
        count += 1

    cap.release()
    return saved_frames

frames = extract_frames(video_path)
print(f"Extracted {len(frames)} frames.")


def detect_slide_changes_ssim(frames, threshold=0.98):
    slides = [frames[0]]  # First frame is always a slide
    prev_img = cv2.imread(frames[0], cv2.IMREAD_GRAYSCALE)

    for i in range(1, len(frames)):
        curr_img = cv2.imread(frames[i], cv2.IMREAD_GRAYSCALE)
        
        if prev_img.shape != curr_img.shape:
            curr_img = cv2.resize(curr_img, (prev_img.shape[1], prev_img.shape[0]))

        score, _ = ssim(prev_img, curr_img, full=True)

        if score < threshold:  # Lower SSIM means more difference
            slides.append(frames[i])
            prev_img = curr_img  # Update the reference slide

    return slides

slides = detect_slide_changes_ssim(frames, threshold=0.98)
print(f"Detected {len(slides)} slides with SSIM method.")


slide_folder = "slides_new"
os.makedirs(slide_folder, exist_ok=True)

for slide in slides:
    slide_name = os.path.basename(slide)
    os.rename(slide, os.path.join(slide_folder, slide_name))

print(f"Slides saved in '{slide_folder}' folder.")


def create_ppt_from_folder(folder_path, output_ppt="slides.pptx"):
    prs = Presentation()

    # Get all image files from the folder and sort them using natural sorting
    slides = natsorted(
        [os.path.join(folder_path, f) for f in os.listdir(folder_path) if f.endswith(('.png', '.jpg', '.jpeg'))]
    )

    if not slides:
        print("No images found in the folder.")
        return

    for slide_img in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide layout
        left = top = Inches(0.5)  # Margins
        pic = slide.shapes.add_picture(slide_img, left, top, width=Inches(9), height=Inches(5))

    prs.save(output_ppt)
    print(f"PPT created: {output_ppt}")

# Set the folder containing images
slides_folder = "slides_new3"
output_ppt = "extracted_slides_4.pptx"
create_ppt_from_folder(slides_folder, output_ppt)






